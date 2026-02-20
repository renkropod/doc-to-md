"""PPTX 파서 - python-pptx 기반"""

import logging
from pathlib import Path

from doc_to_md.exceptions import ParserError
from doc_to_md.parsers.base import BaseParser, ParseResult
from doc_to_md.utils.image_handler import generate_image_filename, save_image

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    """PPTX 문서를 마크다운으로 변환하는 파서"""

    def supported_extensions(self) -> list[str]:
        return [".pptx", ".ppt"]

    def parse(self, file_path: Path, **kwargs) -> ParseResult:
        """PPTX 파일을 슬라이드별 섹션으로 변환"""
        extract_images = kwargs.get("extract_images", False)
        image_dir = kwargs.get("image_dir")

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise ParserError(
                "PPTX 파싱을 위해 python-pptx를 설치해주세요: pip install python-pptx"
            )

        try:
            prs = Presentation(str(file_path))
        except Exception as e:
            raise ParserError(f"PPTX 파일 열기 실패: {e}") from e

        parts: list[str] = []
        tables: list[str] = []
        images: list[dict] = []
        img_index = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_parts: list[str] = [f"## 슬라이드 {slide_num}"]

            # 슬라이드 노트
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"> 노트: {notes}")

            for shape in slide.shapes:
                # 텍스트 프레임
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if not text:
                            continue

                        level = paragraph.level
                        if level == 0 and shape.shape_type == 13:  # Title
                            slide_parts.append(f"### {text}")
                        elif level > 0:
                            indent = "  " * (level - 1)
                            slide_parts.append(f"{indent}- {text}")
                        else:
                            slide_parts.append(text)

                # 테이블
                if shape.has_table:
                    md_table = self._table_to_markdown(shape.table)
                    tables.append(md_table)
                    slide_parts.append(md_table)

                # 이미지
                if extract_images and image_dir and hasattr(shape, "image"):
                    try:
                        img_data = shape.image.blob
                        content_type = shape.image.content_type
                        ext = content_type.split("/")[-1] if content_type else "png"
                        if ext == "jpeg":
                            ext = "jpg"

                        filename = generate_image_filename(
                            file_path.stem, img_index, ext
                        )
                        save_image(img_data, image_dir, filename)
                        rel_path = f"./{image_dir.name}/{filename}"
                        slide_parts.append(f"![Image {img_index}]({rel_path})")
                        images.append({
                            "path": str(image_dir / filename),
                            "alt": f"Image {img_index}",
                            "page": slide_num,
                        })
                        img_index += 1
                    except Exception as e:
                        logger.debug("PPTX 이미지 추출 실패: %s", e)

            parts.append("\n\n".join(slide_parts))

        metadata = self._extract_metadata(prs)

        return ParseResult(
            content="\n\n---\n\n".join(parts),
            metadata=metadata,
            images=images,
            tables=tables,
        )

    def _table_to_markdown(self, table) -> str:
        """PPTX 테이블을 마크다운으로 변환"""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
            rows.append(cells)

        if not rows:
            return ""

        col_count = max(len(r) for r in rows)
        lines = [
            "| " + " | ".join(rows[0] + [""] * (col_count - len(rows[0]))) + " |",
            "| " + " | ".join("---" for _ in range(col_count)) + " |",
        ]

        for row in rows[1:]:
            padded = row + [""] * (col_count - len(row))
            lines.append("| " + " | ".join(padded[:col_count]) + " |")

        return "\n".join(lines)

    def _extract_metadata(self, prs) -> dict:
        """PPTX 메타데이터 추출"""
        metadata = {}
        try:
            props = prs.core_properties
            if props.title:
                metadata["title"] = props.title
            if props.author:
                metadata["author"] = props.author
            if props.subject:
                metadata["subject"] = props.subject
            metadata["slide_count"] = str(len(prs.slides))
        except Exception as e:
            logger.debug("PPTX 메타데이터 추출 실패: %s", e)
        return metadata
